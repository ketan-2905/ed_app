import os
import io
from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image

def extract_pptx_content(file_path):
   
    ppt_content = {
        "text": [],
        "images": [],
        "tables": []
    }
    
    prs = Presentation(file_path)
    output_dir = os.path.dirname(file_path) 
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = " ".join(run.text for run in paragraph.runs)
                    slide_text.append(paragraph_text)
            
            if shape.has_table:
                table = shape.table
                table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                ppt_content["tables"].append(df)
                
            if shape.shape_type == 13: 
                image = shape.image
                image_bytes = io.BytesIO(image.blob)
                ext = image.ext
                image_filename = f"{output_dir}/ppt_slide{slide_num}_img.{ext}"
                
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes.getbuffer())
                
                ppt_content["images"].append(image_filename)
        
        if slide_text:
            ppt_content["text"].append("\n".join(slide_text))
    
    return ppt_content

def preview_extracted_images(images):
    """Preview extracted images using matplotlib."""
    for img_path in images:
        img = Image.open(img_path)
        plt.imshow(img)
        plt.title(os.path.basename(img_path))
        plt.axis('off')
        plt.show()
